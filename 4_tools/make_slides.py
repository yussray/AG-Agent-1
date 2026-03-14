# ============================================================
# make_slides.py — Antigravity Multi-Agent System
# Generate PowerPoint presentations from text outlines
# ============================================================

import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
import paths  # noqa: E402

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import re
from config import OUTPUT_DIR


def parse_outline(outline: str) -> list[dict]:
    """Parse SLIDE N: Title / bullet lines into structured data."""
    slides = []
    current = None

    for line in outline.splitlines():
        line = line.strip()
        if not line:
            continue
        m = re.match(r'^SLIDE\s*(\d+):\s*(.+)', line, re.IGNORECASE)
        if m:
            if current:
                slides.append(current)
            current = {"title": m.group(2).strip(), "bullets": []}
        elif line.startswith("-") and current:
            bullet = line.lstrip("-").strip()
            if bullet:
                current["bullets"].append(bullet)

    if current:
        slides.append(current)
    return slides


def make_slides(outline: str, title: str = "Presentation", filename: str = None) -> str:
    """
    Generate a PPTX file from a text outline.

    Args:
        outline: Slide outline text (SLIDE N: Title / bullets)
        title: Presentation title (used for title slide if no slides parsed)
        filename: Output filename (auto-generated if None)

    Returns:
        Absolute path to saved .pptx file
    """
    slides_data = parse_outline(outline)
    if not slides_data:
        slides_data = [{"title": title, "bullets": ["No content generated. Please try again."]}]

    prs = Presentation()

    # Dark background theme
    BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
    TITLE_COLOR = RGBColor(0xE9, 0x4F, 0x37) # red accent
    BODY_COLOR = RGBColor(0xE0, 0xE0, 0xE0)  # light grey
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BG_COLOR

        # Title box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), SLIDE_W - Inches(1), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR

        # Bullet box
        if slide_data["bullets"]:
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), SLIDE_W - Inches(1), SLIDE_H - Inches(2.5))
            tf2 = body_box.text_frame
            tf2.word_wrap = True

            for i, bullet in enumerate(slide_data["bullets"][:6]):
                p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p2.text = f"• {bullet}"
                p2.font.size = Pt(18)
                p2.font.color.rgb = BODY_COLOR

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    if not filename:
        import re as _re
        slug = _re.sub(r'[^\w]', '_', title.lower())[:25]
        filename = f"slides_{slug}.pptx"

    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    return os.path.abspath(path)
